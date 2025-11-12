#!/usr/bin/env python3
# scripts/latex_to_pptx.py
import sys
import os
from pdf2image import convert_from_path
from pptx import Presentation
from pptx.util import Inches

def pdf_to_pptx(pdf_path, pptx_path, dpi=150):
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF not found: {pdf_path}")
    pages = convert_from_path(pdf_path, dpi=dpi)  # lista de PIL images
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # layout vacío

    for i, img in enumerate(pages):
        slide = prs.slides.add_slide(blank_slide_layout)
        img_path = f"_tmp_slide_{i}.png"
        img.save(img_path, "PNG")
        # insertar imagen ocupando todo el slide
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        os.remove(img_path)

    prs.save(pptx_path)
    print("Saved:", pptx_path)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python scripts/latex_to_pptx.py input.pdf output.pptx")
        sys.exit(1)
    pdf_to_pptx(sys.argv[1], sys.argv[2])
