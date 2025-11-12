#!/usr/bin/env python3
# scripts/text_to_pptx.py
from pptx import Presentation
from pptx.util import Pt

def make_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[1]  # título y contenido
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    # limpia el primer párrafo generado por defecto
    body.clear()
    for b in bullets:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)

def create_pptx_from_structure(structure, out="output.pptx"):
    prs = Presentation()
    for s in structure:
        make_slide(prs, s.get("title",""), s.get("bullets", []))
    prs.save(out)
    print("Saved:", out)

# ejemplo de uso
if __name__ == "__main__":
    sample = [
        {"title":"Slide 1", "bullets":["Punto A", "Punto B"]},
        {"title":"Slide 2", "bullets":["Otro punto"]}
    ]
    create_pptx_from_structure(sample, "example_from_text.pptx")
